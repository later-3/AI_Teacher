from __future__ import annotations

import re
from pathlib import Path
from typing import Iterable, List, Tuple

import pdfplumber
from pptx import Presentation

_HEADING_PATTERN = re.compile(r"^(#{1,6})\s+(.*)")


def parse_pptx(path: Path) -> Iterable[Tuple[int, str]]:
    presentation = Presentation(path)
    for idx, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content = shape.text.strip()
                if content:
                    texts.append(content)
        if texts:
            yield idx, "\n".join(texts)


def parse_pdf(path: Path) -> Iterable[Tuple[int, str]]:
    with pdfplumber.open(path) as pdf:
        for idx, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            text = text.strip()
            if text:
                yield idx, text


def parse_text_file(path: Path) -> Iterable[Tuple[int, str]]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    chunk = []
    idx = 1
    for line in lines:
        chunk.append(line)
        if len(chunk) >= 5:
            yield idx, "\n".join(chunk)
            idx += 1
            chunk = []
    if chunk:
        yield idx, "\n".join(chunk)


def parse_markdown_file(path: Path) -> Iterable[Tuple[int, str]]:
    """Split markdown into sections using headings/code fences awareness."""
    text = path.read_text(encoding="utf-8", errors="ignore")
    lines = text.splitlines()
    idx = 1
    buffer: List[str] = []
    in_code_fence = False

    def flush_buffer() -> Tuple[int, str]:
        nonlocal idx
        content = "\n".join(buffer).strip()
        buffer.clear()
        if content:
            current_idx = idx
            idx += 1
            return current_idx, content
        return 0, ""

    for raw_line in lines:
        line = raw_line.rstrip()
        stripped = line.strip()
        if stripped.startswith("```"):
            in_code_fence = not in_code_fence
            continue
        if not in_code_fence and not stripped:
            continue
        if not in_code_fence:
            heading_match = _HEADING_PATTERN.match(stripped)
            if heading_match:
                section_idx, content = flush_buffer()
                if section_idx:
                    yield section_idx, content
                buffer.append(stripped)
                continue
        buffer.append(raw_line if in_code_fence else stripped)

    if buffer:
        section_idx, content = flush_buffer()
        if section_idx:
            yield section_idx, content
